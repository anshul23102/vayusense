"""Fill the official hackathon PPT template with VayuSense content.

Design system (applied consistently on every content slide):
- Canvas: the template is 10in x 5.625in.
- Font: Google Sans everywhere (the template's own font).
- Cards: light gray rounded rectangles (F8F9FA), no borders, no shadows.
- One accent: Google Blue (1A73E8) for key stats and links;
  a readable green only for the NVIDIA stack heading.
- Sizes: card heading 11.5pt bold, card body 9.5pt, stat numbers 30pt.
- Short copy, no walls of text, no em dashes, no double hyphens.
"""
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

SRC = "docs/deck/template.pptx"
OUT = "docs/deck/VayuSense_Submission_Deck.pptx"

SLIDE_W = 10.0
FONT = "Google Sans"
DARK = RGBColor(0x20, 0x21, 0x24)
BODY = RGBColor(0x3C, 0x40, 0x43)
MUTE = RGBColor(0x5F, 0x63, 0x68)
BLUE = RGBColor(0x1A, 0x73, 0xE8)
NV = RGBColor(0x4C, 0x8C, 0x2B)
CARD = RGBColor(0xF8, 0xF9, 0xFA)

B, I, _ = True, True, False


def shapes_by_id(slide):
    return {sh.shape_id: sh for sh in slide.shapes}


def set_plain(shape, text, size=14, bold=False):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = DARK


def add_runs(p, spec, size, color=BODY):
    for item in spec:
        text, bold, italic = item[0], item[1], item[2]
        col = item[3] if len(item) > 3 else color
        run = p.add_run()
        run.text = text
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = col


def force_font(shape, size=None, space_after=None):
    for p in shape.text_frame.paragraphs:
        if space_after is not None:
            p.space_after = Pt(space_after)
        for r in p.runs:
            r.font.name = FONT
            if size:
                r.font.size = Pt(size)


def card(slide, x, y, w, h):
    sh = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.adjustments[0] = 0.07
    sh.fill.solid()
    sh.fill.fore_color.rgb = CARD
    sh.line.fill.background()
    sh.shadow.inherit = False
    tf = sh.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.13)
    tf.margin_right = Inches(0.13)
    tf.margin_top = Inches(0.09)
    tf.margin_bottom = Inches(0.07)
    return sh


def feature_card(slide, x, y, w, h, title, body, title_color=DARK,
                 title_size=11.5, body_size=9.5):
    sh = card(slide, x, y, w, h)
    tf = sh.text_frame
    p = tf.paragraphs[0]
    p.space_after = Pt(2)
    add_runs(p, [(title, B, _)], title_size, color=title_color)
    p2 = tf.add_paragraph()
    add_runs(p2, body if isinstance(body, list) else [(body, _, _)], body_size)
    return sh


def stat_card(slide, x, y, w, h, number, label):
    sh = card(slide, x, y, w, h)
    tf = sh.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(3)
    add_runs(p, [(number, B, _)], 30, color=BLUE)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    add_runs(p2, [(label, _, _)], 9.5, color=MUTE)
    return sh


def caption(slide, x, y, w, text):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.28))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    add_runs(p, text if isinstance(text, list) else [(text, _, _)], 9.5, color=MUTE)
    return box


def crop_right(src, dst, px):
    im = Image.open(src)
    im.crop((0, 0, im.size[0] - px, im.size[1])).save(dst)


def center_picture(slide, path, y, height=None, width=None):
    im = Image.open(path)
    ratio = im.size[0] / im.size[1]
    if height:
        width = height * ratio
    else:
        height = width / ratio
    x = (SLIDE_W - width) / 2
    return slide.shapes.add_picture(path, Inches(x), Inches(y),
                                    width=Inches(width), height=Inches(height))


crop_right("docs/deck/screenshots/landing.png", "docs/deck/screenshots/landing_c.png", 26)
crop_right("docs/deck/screenshots/dashboard_top.png", "docs/deck/screenshots/dashboard_top_c.png", 26)
crop_right("docs/deck/screenshots/dashboard_trend.png", "docs/deck/screenshots/dashboard_trend_c.png", 26)

prs = Presentation(SRC)
slides = list(prs.slides)

URL = "vayusense-kggf.onrender.com"
X0 = 0.47          # left margin
CW3 = 2.89         # 3-column card width
CW2 = 4.43         # 2-column card width
GAP = 0.19

# ---------------- Slide 1: Participant details ----------------
s = shapes_by_id(slides[0])
tf = s[55].text_frame
for p in tf.paragraphs:
    full = "".join(r.text for r in p.runs)
    if full.strip().startswith("Participant Name"):
        for r in p.runs:
            r.text = ""
            r.font.name = FONT
        p.runs[0].text = "Participant Name: Anshul Jain (Team BloodWyrm, Solo)"
    elif full.strip().startswith("Problem Statement"):
        for r in p.runs:
            r.text = ""
            r.font.name = FONT
        p.runs[0].text = ("Problem Statement: AI for Better Living and Smarter Communities. "
                          "Create a data intelligence tool people would actually use, and show "
                          "how acceleration helps them make a faster or better decision.")

# ---------------- Slide 2: Brief about the idea ----------------
sl = slides[1]
s = shapes_by_id(sl)
set_plain(s[62], " ", size=1)

intro = sl.shapes.add_textbox(Inches(X0), Inches(1.34), Inches(9.06), Inches(0.95))
tf = intro.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.space_after = Pt(4)
add_runs(p, [
    ("VayuSense", B, _), (" is an AI decision intelligence platform that turns millions of raw "
     "air quality readings into one honest answer: ", _, _),
    ("is it safe to go outside right now?", B, I),
], 13, color=BODY)
p2 = tf.add_paragraph()
add_runs(p2, [
    ("Built for parents, schools, clinics, and city officials across Indian and APAC cities.",
     _, _),
], 11, color=MUTE)

stat_card(sl, X0, 2.5, CW3, 1.28, "5.9M+", "real OpenAQ sensor readings, processed on GPU")
stat_card(sl, X0 + CW3 + GAP, 2.5, CW3, 1.28, "37.5×",
          "faster than pandas: NVIDIA RAPIDS cuDF benchmark")
stat_card(sl, X0 + 2 * (CW3 + GAP), 2.5, CW3, 1.28, "2 agents",
          "Google ADK pipeline on Gemini: Analyst → Advisor")

closing = sl.shapes.add_textbox(Inches(X0), Inches(4.12), Inches(9.06), Inches(1.15))
tf = closing.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.space_after = Pt(4)
add_runs(p, [
    ("It goes beyond a raw AQI number: a ", _, _), ("Safety Score", B, _), (", a ", _, _),
    ("Human Impact", B, _), (" estimate (cigarette-equivalent exposure, life-expectancy impact), "
     "and a plain-language verdict anyone can act on.", _, _),
], 11, color=BODY)
p2 = tf.add_paragraph()
add_runs(p2, [
    ("Fully deployed and live at ", _, _), (URL, B, _, BLUE), (".", _, _),
], 11, color=BODY)

# ---------------- Slide 3: Approach / impact / architecture ----------------
sl = slides[2]
s = shapes_by_id(sl)
force_font(s[70], size=10, space_after=3)
set_plain(s[68], " ", size=1)

y0, ch = 3.42, 1.6
feature_card(sl, X0, y0, CW3, ch, "Approach", [
    ("Ingest ", _, _), ("5.9M+ OpenAQ readings", B, _),
    ("; clean and trend-analyze on GPU with ", _, _), ("NVIDIA cuDF", B, _),
    ("; reason with ", _, _), ("Google ADK agents", B, _),
    (" (Analyst → Advisor) on ", _, _), ("Gemini", B, _),
    ("; ship as a deployed FastAPI app.", _, _),
], title_color=BLUE, body_size=9)
feature_card(sl, X0 + CW3 + GAP, y0, CW3, ch, "Real-world impact", [
    ("Delhi NCR breathes hazardous air most days. VayuSense turns raw data into ", _, _),
    ("same-day decisions", B, _),
    (" for schools, clinics, and city wards; not another chart nobody reads.", _, _),
], title_color=BLUE, body_size=9)
feature_card(sl, X0 + 2 * (CW3 + GAP), y0, CW3, ch, "Architecture", [
    ("OpenAQ → NVIDIA cuDF (GPU) → ADK/Gemini agents → FastAPI dashboard. Raw readings become a ",
     _, _), ("Safety Score", B, _), (", a ", _, _),
    ("Human Impact", B, _), (" estimate, and a clear recommendation.", _, _),
], title_color=BLUE, body_size=9)

# ---------------- Slide 4: Opportunities / USP ----------------
sl = slides[3]
s = shapes_by_id(sl)
set_plain(s[77], " ", size=1)

r1y, r1h = 2.32, 1.18
feature_card(sl, X0, r1y, CW3, r1h, "Proven acceleration, not a claim",
             [("A reproducible ", _, _), ("37.5× cuDF vs pandas benchmark", B, _),
              (" over 5.9M+ real rows, shown live inside the app.", _, _)], body_size=9)
feature_card(sl, X0 + CW3 + GAP, r1y, CW3, r1h, "Human stakes, not just AQI",
             [("Cigarette-equivalent exposure and life-expectancy impact, grounded in "
               "published epidemiology.", _, _)], body_size=9)
feature_card(sl, X0 + 2 * (CW3 + GAP), r1y, CW3, r1h, "Agents, not a prompt wrapper",
             [("ADK Sequential agents: the ", _, _), ("Analyst", B, _),
              (" gathers facts with real tools; the ", _, _), ("Advisor", B, _),
              (" turns them into guidance.", _, _)], body_size=9)
r2y, r2h = 3.68, 1.05
w2 = 4.43
x1 = (SLIDE_W - (2 * w2 + GAP)) / 2
feature_card(sl, x1, r2y, w2, r2h, "Answers the brief end to end",
             [("A real ingest, clean, model, visualize pipeline; 2+ Google Cloud items; "
               "2+ NVIDIA items; a genuinely useful output.", _, _)], body_size=9)
feature_card(sl, x1 + w2 + GAP, r2y, w2, r2h, "Built to scale beyond the hackathon",
             [("The same architecture serves any city or pollutant dataset; "
               "just repoint the ingestion.", _, _)], body_size=9)

# ---------------- Slide 5: Features ----------------
sl = slides[4]
s = shapes_by_id(sl)
set_plain(s[84], " ", size=1)

y0 = 1.42
ch, vgap = 0.84, 0.12
feats = [
    ("Live multi-city dashboard",
     [("Delhi, Mumbai, Kolkata, Chennai: 90-day trends, 7-day rolling averages", _, _)]),
    ("Safety Score (0 to 100)",
     [("one plain-language verdict per city, driven by WHO exceedance", _, _)]),
    ("Human Impact panel",
     [("cigarette-equivalent exposure and life-expectancy impact", _, _)]),
    ("Pollution hotspot ranking",
     [("the worst stations in each city, ranked by average PM2.5", _, _)]),
    ("Ask VayuSense chat",
     [("ADK multi-agent pipeline (Analyst → Advisor) on Gemini", _, _)]),
    ("GPU-accelerated analytics",
     [("NVIDIA RAPIDS cuDF backend with a live 37.5× benchmark panel", _, _)]),
    ("WHO-guideline comparisons",
     [("six pollutants tracked against WHO 24-hour limits, with trends", _, _)]),
    ("Deployed and public",
     [("live right now at ", _, _), (URL, B, _, BLUE)]),
]
for i, (t, b) in enumerate(feats):
    row, col = divmod(i, 2)
    feature_card(sl, X0 + col * (CW2 + GAP), y0 + row * (ch + vgap), CW2, ch, t, b,
                 title_size=10.5, body_size=9)

# ---------------- Slide 6: Process flow diagram ----------------
s = shapes_by_id(slides[5])
set_plain(s[91], " ", size=1)
center_picture(slides[5], "docs/deck/diagram_flow.png", 1.32, width=8.55)
caption(slides[5], 1.0, 4.95, 8.0,
        "One question in, one decision out: the ADK pipeline grounds every answer in GPU-processed data.")

# ---------------- Slide 7: Wireframe / landing page ----------------
s = shapes_by_id(slides[6])
set_plain(s[98], " ", size=1)
center_picture(slides[6], "docs/deck/screenshots/landing_c.png", 1.38, height=3.45)
caption(slides[6], 1.0, 4.98, 8.0,
        [("Landing page of the deployed product: ", _, _), (URL, B, _, BLUE)])

# ---------------- Slide 8: Architecture diagram ----------------
s = shapes_by_id(slides[7])
set_plain(s[105], " ", size=1)
center_picture(slides[7], "docs/deck/diagram_arch.png", 1.42, height=3.95)

# ---------------- Slide 9: Tech stack ----------------
sl = slides[8]
s = shapes_by_id(sl)
force_font(s[111])
set_plain(s[112], " ", size=1)

y0, ch = 1.78, 1.28
sh = card(sl, X0, y0, CW2, ch)
tf = sh.text_frame
p = tf.paragraphs[0]
p.space_after = Pt(3)
add_runs(p, [("Google Cloud", B, _)], 11.5, color=BLUE)
for line in [
    [("Gemini 2.5 Flash", B, _), (" powering both agents' reasoning", _, _)],
    [("Google ADK", B, _), (" Sequential multi-agent orchestration, tool calling", _, _)],
    [("Cloud Run-ready", B, _), (" containerized FastAPI app (Docker)", _, _)],
]:
    pp = tf.add_paragraph()
    pp.space_after = Pt(3)
    add_runs(pp, line, 9.5)

sh = card(sl, X0 + CW2 + GAP, y0, CW2, ch)
tf = sh.text_frame
p = tf.paragraphs[0]
p.space_after = Pt(3)
add_runs(p, [("NVIDIA", B, _)], 11.5, color=NV)
for line in [
    [("RAPIDS cuDF", B, _), (" GPU dataframes: clean, aggregate, trend, anomaly", _, _)],
    [("T4 GPU", B, _), (" (Google Colab) as the benchmark target", _, _)],
    [("9.34s → 0.25s", B, _, BLUE), (" on 5.9M rows: a 37.5× measured speedup", _, _)],
]:
    pp = tf.add_paragraph()
    pp.space_after = Pt(3)
    add_runs(pp, line, 9.5)

sh = card(sl, X0, 3.28, 2 * CW2 + GAP, 1.42)
tf = sh.text_frame
p = tf.paragraphs[0]
p.space_after = Pt(3)
add_runs(p, [("Why this stack", B, _)], 11.5, color=DARK)
for line in [
    [("Speed: ", B, _), ("millions of sensor rows need cuDF-class processing to stay "
      "decision-fresh instead of report-stale.", _, _)],
    [("Trust: ", B, _), ("grounded, tool-calling ADK agents avoid the hallucinated health "
      "claims a single freeform prompt risks.", _, _)],
    [("Portability: ", B, _), ("one lightweight container deploys to Cloud Run, Render, or "
      "on-prem with no code changes.", _, _)],
]:
    pp = tf.add_paragraph()
    pp.space_after = Pt(3)
    add_runs(pp, line, 9.5)

# ---------------- Slide 10: Snapshots of the running prototype ----------------
sl = slides[9]
s = shapes_by_id(sl)
set_plain(s[120], " ", size=1)
iw = 4.43
im = Image.open("docs/deck/screenshots/dashboard_top_c.png")
ih = iw / (im.size[0] / im.size[1])
y0 = 1.38
sl.shapes.add_picture("docs/deck/screenshots/dashboard_top_c.png",
                      Inches(X0), Inches(y0), width=Inches(iw))
sl.shapes.add_picture("docs/deck/screenshots/dashboard_trend_c.png",
                      Inches(X0 + iw + GAP), Inches(y0), width=Inches(iw))
caption(sl, X0, y0 + ih + 0.1, iw,
        "Safety Score, WHO comparisons and Human Impact, live for Delhi")
caption(sl, X0 + iw + GAP, y0 + ih + 0.1, iw,
        "90-day trends, the live GPU benchmark and the agent chat")
final = sl.shapes.add_textbox(Inches(X0), Inches(y0 + ih + 0.52), Inches(9.06), Inches(0.35))
tf = final.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
add_runs(p, [("Try it live: ", _, _), (URL, B, _, BLUE)], 11, color=BODY)

prs.save(OUT)
print("Saved", OUT)
